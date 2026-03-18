"""
ah_slides.py,Asset Hub Analysis Slide Deck
=============================================
Generates a PowerPoint presentation from the analysis outputs.
Uses the pre-rendered PNGs and analysis results from the DuckDB database.

Usage:
    python ah_slides.py              # uses latest run
    python ah_slides.py --run-id 1   # specific run
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

import duckdb
from ah_loader import DB_FILE

OUT_DIR = Path("outputs")

# ── AWS Palette ───────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x23, 0x2F, 0x3E)   # Squid Ink
BG_SLIDE   = RGBColor(0xFA, 0xFA, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x54, 0x60, 0x6E)   # AWS Slate
GREEN      = RGBColor(0x1B, 0x66, 0x0F)   # AWS Green (success)
BLUE       = RGBColor(0x00, 0x73, 0xBB)   # AWS Anchor Blue
RED        = RGBColor(0xD1, 0x34, 0x12)   # AWS Red (error/critical)
AMBER      = RGBColor(0xFF, 0x99, 0x00)   # AWS Orange (primary brand)
LIGHT_GRAY = RGBColor(0xEB, 0xED, 0xF0)   # AWS Cloud Gray
FONT       = 'Amazon Ember'


def _set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *,
              font_size=18, bold=False, color=BG_DARK, alignment=PP_ALIGN.LEFT,
              font_name=None):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or FONT
    return tf


def _add_bullet_list(slide, left, top, width, height, items, *,
                     font_size=14, color=BG_DARK):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = FONT
    return tf


def _add_table(slide, left, top, width, row_height, headers, rows, *,
               header_color=GREEN):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_width = width // n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, row_height * n_rows
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = FONT

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = BG_DARK
                    run.font.name = FONT

    return table_shape


def _add_link_button(slide, left, top, width, height, text, url, *,
                     fill_color=AMBER):
    """Add a clickable button shape that opens a URL."""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT

    # Add hyperlink via oxml
    rPr = run._r.get_or_add_rPr()
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    hlinkClick.set(qn('r:id'), '')
    rPr.append(hlinkClick)
    # Set the hyperlink relationship
    rel = slide.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hlinkClick.set(qn('r:id'), rel)


def _slide_title(slide, title, subtitle=None):
    """Standard slide title bar."""
    _set_slide_bg(slide, BG_SLIDE)

    # Title band
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(0.9)  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.fill.background()

    # AWS orange accent line
    accent = slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.88), Inches(12.1), Inches(0.03)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    _add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5),
              title, font_size=24, bold=True, color=WHITE)

    if subtitle:
        _add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.35),
                  subtitle, font_size=12, color=GRAY)


def generate_slides(run_id: int | None = None, out_path: str = 'outputs/slides.pptx'):
    """Generate the full slide deck."""
    # ── Load data ────────────────────────────────────────────────────────────
    con = duckdb.connect(DB_FILE, read_only=True)

    if run_id is None:
        run_id = con.execute("SELECT MAX(run_id) FROM runs").fetchone()[0]

    run_ts = str(con.execute(
        "SELECT timestamp FROM runs WHERE run_id = ?", [run_id]
    ).fetchone()[0])[:19]

    n_concepts = con.execute(
        "SELECT COUNT(DISTINCT concept) FROM concepts WHERE run_id = ?", [run_id]
    ).fetchone()[0]

    n_rels = con.execute(
        "SELECT COUNT(*) FROM relationships WHERE run_id = ?", [run_id]
    ).fetchone()[0]

    con.close()

    # ── Create presentation (widescreen 16:9) ───────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 1,Title
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, BG_DARK)

    # AWS orange accent bar
    bar = slide.shapes.add_shape(
        1, Inches(1), Inches(1.5), Inches(1.5), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    _add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
              'Asset Hub', font_size=44, bold=True, color=WHITE)
    _add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.7),
              'Data Structure Analysis', font_size=28, color=GRAY)
    _add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
              f'{n_concepts} concepts  \u00b7  {n_rels} relationships  \u00b7  '
              f'Run #{run_id}  \u00b7  {run_ts}',
              font_size=14, color=GRAY)
    _add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
              '294 of 350 relationships (84%) cannot be resolved today.\n'
              'Trangis provides 0 GUIDs out of 29 concepts.',
              font_size=18, bold=True, color=AMBER)

    slide.notes_slide.notes_text_frame.text = (
        "We analysed all 76 concepts and 350 relationships in Asset Hub. "
        "The numbers on this slide are the headline: 84% of relationships are unresolved, "
        "and the root cause traces back to one source system providing zero GUIDs. "
        "The rest of this deck shows exactly where those failures are, why they exist, "
        "and what it would take to fix them."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 2,Context
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Asset Hub: Two Source Systems, Two Realities')

    items = [
        'Asset Hub is a key-value store: every concept must have a unique, stable identifier (GUID)',
        'InfraRef (MSSQL): 47 concepts, 44 with GUID (93.6%)',
        'Trangis (Oracle views): 29 concepts, 0 with GUID (0.0%)',
        'Trangis views are read-only, the schema cannot be altered, '
        'the system is scheduled for decommissioning',
        'Missing GUIDs: 32 concepts that cannot be addressed in the store',
        '294 relationships that cannot be traversed because the target has no GUID',
    ]
    _add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
                     items, font_size=16)

    slide.notes_slide.notes_text_frame.text = (
        "These are the raw numbers. InfraRef provides GUIDs for 44 of 47 concepts, "
        "Trangis provides GUIDs for 0 of 29. The 3 InfraRef exceptions are legacy tables "
        "that can be remediated. The 29 Trangis gaps cannot, because the Oracle views are "
        "read-only and the system is in decommissioning. This distinction matters: "
        "one side is a data quality issue with a fix path, the other is an architectural "
        "constraint with no fix path within the current system."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 3,Concept Relationship Graph
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Concept Relationship Graph',
                 'Node size = in-degree  \u00b7  Green = InfraRef  \u00b7  '
                 'Red = Trangis  \u00b7  Yellow border = no GUID or non-unique ID')

    graph_path = str(OUT_DIR / 'graph.png')
    slide.shapes.add_picture(
        graph_path, Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.3)
    )

    _add_link_button(slide, Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35),
                     'Open interactive graph', 'graph.html')

    slide.notes_slide.notes_text_frame.text = (
        "This is a direct rendering of the data, not a model. Every dot is a concept, "
        "every line is a relationship stored in the database. Green nodes are InfraRef, "
        "red nodes are Trangis, turquoise is GeographicalSite which exists in both. "
        "Yellow borders indicate a data quality issue: no GUID or non-unique ID. "
        "Note that the entire Trangis cluster is yellow-bordered. "
        "The InfraRef side has only a few yellow nodes (3 legacy exceptions)."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 4,Graph Interpretation
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Graph Topology: Hub Nodes and Resolution Rates')

    left_items = [
        '6 hub nodes receive the majority of all incoming edges',
        'ServiceCenter: 53 incoming references, 1.9% GUID resolution',
        'WorkCenter: 44 incoming references, 100% GUID resolution',
        'GeographicalBay (40), PowerTransformerPlace (40), '
        'GeographicalSubstation (42), GeographicalContainer (30)',
        'WorkCenter demonstrates the architecture works when GUIDs are present',
    ]
    _add_text(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
              'Hub nodes (largest dots)', font_size=16, bold=True, color=GREEN)
    _add_bullet_list(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.8),
                     left_items, font_size=13)

    right_items = [
        '29 of 29 Trangis nodes have yellow borders (0% GUID coverage)',
        'GeographicalSite (turquoise): only concept in both systems, '
        'no designated master, join on GeographicalSiteID = P5COD',
        'Cross-boundary edges (InfraRef to Trangis targets) cannot resolve '
        'because the target system provides no GUIDs',
        'Trangis has its own internal topology (Tower, Junction, Line, '
        'GeographicalSpan) but it is entirely disconnected from the GUID graph',
    ]
    _add_text(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4),
              'Trangis boundary', font_size=16, bold=True, color=RED)
    _add_bullet_list(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.8),
                     right_items, font_size=13)

    _add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
              'Implication', font_size=16, bold=True, color=AMBER)
    _add_bullet_list(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
                     [
                         'The 84% unresolved rate is not distributed evenly: InfraRef-to-InfraRef '
                         'relationships resolve well, Trangis relationships resolve at 0%',
                         'Improving the overall rate requires addressing Trangis at source, '
                         'not patching Asset Hub',
                     ], font_size=14, color=BG_DARK)

    slide.notes_slide.notes_text_frame.text = (
        "The graph reveals a star topology, almost everything converges on 5-6 hub nodes. "
        "ServiceCenter is the single most connected node at 53 incoming references, used by "
        "both InfraRef and Trangis. WorkCenter is the success story, 44 incoming references "
        "and 100% GUID resolution, proving the pattern works when the source system cooperates. "
        "The ~40 small green dots in the center are InfraRef equipment assets, each one has "
        "5-7 outgoing edges to the same hubs, creating the visible fan pattern. "
        "On the Trangis side, the red cluster has its own internal topology,Tower, Junction, "
        "GeographicalSpan, Line, but every cross-boundary edge to InfraRef is broken because "
        "GUIDs cannot travel across the Trangis boundary. "
        "The key takeaway: the architecture depends on hub nodes for resolution, but 84% of "
        "the wiring to reach them is broken. This is structural, not a data quality issue."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 5,InfraRef Graph
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'InfraRef, Concept Graph',
                 '47 concepts  \u00b7  Green = InfraRef  \u00b7  '
                 'Red nodes = Trangis targets referenced across the boundary')

    ir_path = str(OUT_DIR / 'graph_infraref.png')
    slide.shapes.add_picture(
        ir_path, Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.3)
    )

    _add_link_button(slide, Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35),
                     'Open interactive graph', 'graph_infraref.html')

    slide.notes_slide.notes_text_frame.text = (
        "This is the InfraRef subgraph in isolation. The fan pattern is very clear here,"
        "~40 equipment assets all pointing to the same 5-6 location hubs. Every equipment "
        "concept follows the same template: references to WorkCenter, ServiceCenter, "
        "GeographicalSubstation, GeographicalBay, PowerTransformerPlace, and often "
        "GeographicalContainer. WorkCenter achieves 100% GUID resolution. The red nodes "
        "you see are Trangis concepts that InfraRef references across the system boundary,"
        "these are the edges where GUIDs get lost. The three InfraRef exceptions without "
        "GUIDs,FaultRecorderAsset, SubStation, VoltageLevel, are visible with yellow borders."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 6,Trangis Graph
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Trangis, Concept Graph',
                 '29 concepts  \u00b7  Red = Trangis  \u00b7  '
                 'Turquoise = dual-source (GeographicalSite)')

    tr_path = str(OUT_DIR / 'graph_trangis.png')
    slide.shapes.add_picture(
        tr_path, Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.3)
    )

    _add_link_button(slide, Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35),
                     'Open interactive graph', 'graph_trangis.html')

    slide.notes_slide.notes_text_frame.text = (
        "The Trangis subgraph has a very different structure from InfraRef. Instead of a "
        "uniform fan pattern, you see an interconnected network topology: Tower, Junction, "
        "GeographicalSpan, and Line form a densely connected core representing the physical "
        "transmission network. Nearly every node has a yellow border, zero GUID coverage. "
        "GeographicalSite in turquoise is the bridge to InfraRef, but without GUIDs on the "
        "Trangis side, this bridge cannot carry identity. The green nodes you see are InfraRef "
        "concepts that Trangis references,ServiceCenter, WorkCenter, GeographicalSubstation,"
        "the relationship exists but cannot be resolved. Line and GuardCircuit have the "
        "additional problem of non-unique IDs from the sign-flipping convention."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 7,GUID Coverage
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'GUID Coverage: 93.6% vs 0.0%')

    _add_table(slide, Inches(0.6), Inches(1.3), Inches(7), Inches(0.45),
               ['System', 'Total', 'With GUID', 'Without GUID', 'Coverage'],
               [
                   ['InfraRef', '47', '44', '3', '93.6%'],
                   ['Trangis', '29', '0', '29', '0.0%'],
                   ['Total', '76', '44', '32', '57.9%'],
               ])

    _add_text(slide, Inches(0.6), Inches(3.2), Inches(6), Inches(0.4),
              'InfraRef exceptions (3 concepts without GUID):',
              font_size=14, bold=True)
    _add_bullet_list(slide, Inches(0.6), Inches(3.6), Inches(6), Inches(1.5),
                     ['FaultRecorderAsset, legacy FaultAnalysis table',
                      'SubStation, FunctionalSite table',
                      'VoltageLevel, FunctionalStation table'],
                     font_size=13)

    _add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1),
              'Trangis: all 29 concepts exposed via read-only Oracle views (VIEW_*). '
              'Schema cannot be altered. System scheduled for decommissioning.',
              font_size=14, color=GRAY)

    slide.notes_slide.notes_text_frame.text = (
        "The GUID divide is stark: InfraRef is at 94%, Trangis is at zero. "
        "InfraRef's 3 exceptions are legacy tables,FaultAnalysis, FunctionalSite, "
        "FunctionalStation, that predate the GUID standard. These are fixable. "
        "Trangis is a different story entirely. These are Oracle views, read-only, "
        "schema cannot be altered, and the system is scheduled for decommissioning. "
        "There is no technical path to adding GUIDs to Trangis short of migrating "
        "the data to a new system. The taxonomy chart on the right shows our five "
        "severity categories, from silent data loss risk at the top to cascading "
        "resolution failures at the bottom."
    )

    # Taxonomy graphic on the right
    tax_path = str(OUT_DIR / 'graphic1_taxonomy.png')
    slide.shapes.add_picture(
        tax_path, Inches(7.8), Inches(1.2), Inches(5.2), Inches(5.0)
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 6,Key Violations
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Non-Unique IDs and Dual-Source Mastery: 4 + 1 Violations')

    _add_text(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
              'Non-Unique IDs (4 concepts)', font_size=18, bold=True, color=RED)

    _add_table(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.4),
               ['Concept', 'System', 'Root Cause'],
               [
                   ['GeographicalBay', 'InfraRef', 'Data quality'],
                   ['PowerTransformerPlace', 'InfraRef', 'Data quality'],
                   ['Line', 'Trangis', 'Sign-flipping'],
                   ['GuardCircuit', 'Trangis', 'Sign-flip cascade'],
               ], header_color=RED)

    _add_text(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(2.5),
              'Line sign-flipping: a positive ID (12.110) means under tension, '
              'negative (\u221212.110) means tension removed. Same physical line, '
              'two keys in the store. GuardCircuit references Line and inherits '
              'the same ambiguity.',
              font_size=13, color=GRAY)

    _add_text(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.4),
              'Dual-Source Mastery', font_size=18, bold=True, color=AMBER)

    _add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4),
                     [
                         'GeographicalSite exists in BOTH InfraRef and Trangis',
                         'Join: InfraRef.GeographicalSiteID = Trangis.P5COD',
                         'Full outer join, a site can exist in only one system',
                         'No designated master, conflicting data is unresolved',
                         'Parcel (Trangis) references GeographicalSite but '
                         'cannot resolve via GUID \u2192 orphaned records',
                     ], font_size=14)

    slide.notes_slide.notes_text_frame.text = (
        "Four concepts have non-unique IDs, this is the most dangerous violation because "
        "a key-value store will silently overwrite data when two records share the same key. "
        "The InfraRef ones,GeographicalBay and PowerTransformerPlace, are data quality "
        "issues that can be investigated and fixed. The Trangis ones are more interesting. "
        "Line uses a sign-flipping convention: positive ID means the line is under tension, "
        "negative means tension removed. So the same physical line appears as two keys,"
        "12.110 and -12.110. GuardCircuit references Line and inherits the same problem. "
        "The cascade impact is an open question with the source team. "
        "On the right: GeographicalSite is the only concept mastered in both systems. "
        "The join on GeographicalSiteID = P5COD is a full outer join, so a site can exist "
        "in only one system. There is no designated master, so conflicting data has no "
        "resolution rule. Parcel from Trangis references GeographicalSite but cannot "
        "resolve the link via GUID, those records become orphans."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 7,Relationship Resolution
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, '294 of 350 Relationships Unresolved (84%)')

    res_path = str(OUT_DIR / 'graphic3_resolution.png')
    slide.shapes.add_picture(
        res_path, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5)
    )

    _add_text(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
              '350 relationships  \u00b7  45 GUID in source  \u00b7  '
              '11 fetched by Asset Hub  \u00b7  294 unresolved (84%)',
              font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    slide.notes_slide.notes_text_frame.text = (
        "This is the headline number: 84% of relationships are unresolved. "
        "Out of 350 relationship edges, only 45 have a GUID available directly in the "
        "source data, and Asset Hub fetches another 11 itself, that's 56 total resolved. "
        "The remaining 294 are broken links. The bar chart on the left shows resolution "
        "rate per target concept. WorkCenter stands out at 100%, proof that the system "
        "works when GUIDs are available. ServiceCenter, despite being referenced 53 times, "
        "resolves under 2%. The donut on the right summarises the overall state. "
        "To-one relationships resolve at 15%, to-many at 21%, both critically low. "
        "This means you cannot traverse the data graph. If you start at an equipment asset "
        "and try to follow its relationships to find its location, service center, or "
        "work center, you will hit a dead end 84% of the time."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 8,Recommendations
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_title(slide, 'Recommendations')

    _add_text(slide, Inches(0.6), Inches(1.2), Inches(3.5), Inches(0.4),
              'Immediate (InfraRef fixes)', font_size=18, bold=True, color=RED)
    _add_bullet_list(slide, Inches(0.6), Inches(1.7), Inches(3.5), Inches(2.5),
                     [
                         'Resolve non-unique IDs for GeographicalBay '
                         'and PowerTransformerPlace',
                         'Add GUIDs to FaultRecorderAsset, SubStation, '
                         'VoltageLevel (3 legacy gaps)',
                         'Confirm cascade behaviour of Trangis line '
                         'sign-flipping with source analysts',
                     ], font_size=13)

    _add_text(slide, Inches(4.6), Inches(1.2), Inches(3.8), Inches(0.4),
              'Medium-term (containment)', font_size=18, bold=True, color=AMBER)
    _add_bullet_list(slide, Inches(4.6), Inches(1.7), Inches(3.8), Inches(2.5),
                     [
                         'Define interim GUID mapping for Trangis '
                         'concepts using business keys',
                         'Establish conflict resolution rule for '
                         'GeographicalSite dual-source merge',
                         'Implement resolution rate monitoring '
                         'to track progress against the 84% baseline',
                     ], font_size=13)

    _add_text(slide, Inches(9), Inches(1.2), Inches(3.8), Inches(0.4),
              'Strategic (Trangis exit)', font_size=18, bold=True, color=GREEN)
    _add_bullet_list(slide, Inches(9), Inches(1.7), Inches(3.8), Inches(2.5),
                     [
                         'Prioritise Trangis successor with GUID '
                         'discipline matching InfraRef',
                         'Migrate critical Trangis concepts (Line, '
                         'Junction, Tower, GeographicalSpan) first',
                         'Design Asset Hub v2 ingestion with '
                         'referential integrity checks at source',
                     ], font_size=13)

    _add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2),
              'The InfraRef side can reach near-100% resolution with targeted fixes. '
              'The Trangis side cannot improve without a successor system. '
              'WorkCenter (100% resolution, 44 references) demonstrates that the Asset Hub '
              'architecture works when the source system provides GUIDs.',
              font_size=16, bold=True, color=BG_DARK)

    slide.notes_slide.notes_text_frame.text = (
        "Three columns, three time horizons. The immediate fixes are all on the InfraRef side "
        "and are achievable: 2 non-unique ID issues, 3 missing GUIDs, and one open question "
        "about line sign-flipping that needs an answer from the Trangis source analysts. "
        "Medium-term is about containing the Trangis problem: interim GUID mappings using "
        "business keys, a merge rule for GeographicalSite, and monitoring so we can track "
        "whether the 84% baseline improves. "
        "Strategic is the real conversation: Trangis needs a successor, and that successor "
        "needs to provide GUIDs. WorkCenter proves the architecture works at 100% when the "
        "source cooperates. The question for stakeholders is not whether to replace Trangis, "
        "but how quickly, and whether Asset Hub v2 should enforce GUID availability at ingestion."
    )

    # ── Save ─────────────────────────────────────────────────────────────────
    prs.save(out_path)
    print(f"  \u2713  {out_path}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate Asset Hub slide deck')
    parser.add_argument('--run-id', type=int, default=None)
    args = parser.parse_args()
    generate_slides(run_id=args.run_id)
